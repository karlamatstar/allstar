"""교차검증 A~D 결과를 3장 발표용 PPT로 변환한다."""

from __future__ import annotations

import csv
import statistics
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from allstar.shared.paths import VOC_REPORT_ROOT

ROOT = Path(__file__).resolve().parent
REPORT_ROOT = VOC_REPORT_ROOT / "cross_validation"
OUTPUT_PATH = REPORT_ROOT / "교차검증_발표자료.pptx"

FONT = "맑은 고딕"
NAVY = RGBColor(20, 31, 55)
BLUE = RGBColor(47, 111, 237)
SKY = RGBColor(84, 166, 255)
GOLD = RGBColor(243, 177, 57)
MINT = RGBColor(40, 185, 150)
RED = RGBColor(224, 91, 91)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(243, 247, 252)
MID = RGBColor(101, 116, 139)
INK = RGBColor(29, 42, 68)
LINE = RGBColor(215, 224, 236)


@dataclass(frozen=True)
class Experiment:
    name: str
    generation: str
    judge: str
    average: float
    median_score: float
    median_seconds: float
    scored: int
    na_count: int
    exception_pass: int


def _numeric(value: str) -> bool:
    try:
        float(value)
        return True
    except (TypeError, ValueError):
        return False


def load_results() -> list[Experiment]:
    providers = {
        "A": ("OpenAI", "Anthropic"),
        "B": ("Anthropic", "OpenAI"),
        "C": ("OpenAI", "OpenAI"),
        "D": ("Anthropic", "Anthropic"),
    }
    results: list[Experiment] = []
    for name, (generation, judge) in providers.items():
        csv_path = REPORT_ROOT / name.lower() / "llm_judge_result.csv"
        with csv_path.open("r", encoding="utf-8-sig", newline="") as handle:
            rows = list(csv.DictReader(handle))
        if len(rows) != 20:
            raise RuntimeError(f"{name} 실험군 결과가 20건이 아닙니다: {len(rows)}건")
        scored_rows = [row for row in rows if _numeric(row.get("total", ""))]
        scores = [float(row["total"]) for row in scored_rows]
        elapsed = [
            float(row["total_seconds"])
            for row in rows
            if row.get("mode") == "live" and _numeric(row.get("total_seconds", ""))
        ]
        results.append(
            Experiment(
                name=name,
                generation=generation,
                judge=judge,
                average=statistics.mean(scores),
                median_score=statistics.median(scores),
                median_seconds=statistics.median(elapsed),
                scored=len(scored_rows),
                na_count=sum(row.get("total") == "N/A" for row in rows),
                exception_pass=sum(row.get("verdict") == "PASS (예외처리)" for row in rows),
            )
        )
    return results


def add_text(slide, x, y, w, h, text, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, radius=True, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_header(slide, number: str, title: str, subtitle: str):
    add_rect(slide, 0, 0, 13.333, 0.13, BLUE, radius=False)
    add_text(slide, 0.55, 0.30, 0.52, 0.46, number, 17, WHITE, True, PP_ALIGN.CENTER)
    badge = slide.shapes[-1]
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.color.rgb = BLUE
    add_text(slide, 1.20, 0.27, 11.5, 0.45, title, 25, NAVY, True)
    add_text(slide, 1.20, 0.72, 11.5, 0.34, subtitle, 10.5, MID)


def style_cell(cell, text, size=11, color=INK, bold=False, fill=WHITE,
               align=PP_ALIGN.CENTER):
    cell.text = str(text)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_summary(prs: Presentation, results: list[Experiment]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    add_header(slide, "01", "A·B·C·D 교차검증 결과", "20개 테스트 케이스 × 4개 모델 조합 · 총 80건")

    table_shape = slide.shapes.add_table(5, 7, Inches(0.55), Inches(1.35), Inches(12.23), Inches(3.15))
    table = table_shape.table
    widths = [0.65, 1.55, 1.55, 1.15, 1.15, 1.35, 1.1]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    headers = ["군", "생성 모델", "평가 모델", "평균점수", "중앙값", "중앙시간", "N/A"]
    for col, header in enumerate(headers):
        style_cell(table.cell(0, col), header, 11, WHITE, True, NAVY)
    colors = [BLUE, GOLD, SKY, MINT]
    for row_idx, (item, accent) in enumerate(zip(results, colors), start=1):
        values = [item.name, item.generation, item.judge, f"{item.average:.1f}",
                  f"{item.median_score:.1f}", f"{item.median_seconds:.2f}초", item.na_count]
        for col, value in enumerate(values):
            fill = RGBColor(235, 242, 255) if row_idx % 2 else WHITE
            style_cell(table.cell(row_idx, col), value, 11.5,
                       accent if col == 0 else INK, col in (0, 3), fill)

    cards = [
        ("최고 평균점수", max(results, key=lambda x: x.average), BLUE),
        ("가장 빠른 중앙시간", min(results, key=lambda x: x.median_seconds), MINT),
        ("API 실패", None, GOLD),
    ]
    x_positions = [0.55, 4.60, 8.65]
    for x, (label, item, accent) in zip(x_positions, cards):
        add_rect(slide, x, 4.88, 3.75, 1.35, WHITE, line=LINE)
        add_rect(slide, x, 4.88, 0.09, 1.35, accent, radius=False)
        add_text(slide, x + 0.25, 5.02, 3.2, 0.28, label, 10.5, MID, True)
        if item:
            value = f"{item.name}  {item.average:.1f}점" if "점수" in label else f"{item.name}  {item.median_seconds:.2f}초"
        else:
            value = "0건 · N/A 없음"
        add_text(slide, x + 0.25, 5.33, 3.2, 0.52, value, 20, NAVY, True)
    add_text(slide, 0.62, 6.58, 12.0, 0.35,
             "정식 점수는 TC-01~16, 데이터 없음 예외처리는 TC-17~18, 장애 재현은 TC-19~20에서 확인",
             10.5, MID, False, PP_ALIGN.CENTER)


def _style_chart(chart, maximum, major_unit, number_format):
    chart.has_legend = False
    chart.has_title = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = maximum
    chart.value_axis.major_unit = major_unit
    chart.value_axis.tick_labels.font.name = FONT
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.number_format = number_format
    chart.category_axis.tick_labels.font.name = FONT
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.category_axis.tick_labels.font.bold = True
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.name = FONT
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.bold = True
    plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.number_format = number_format
    chart.chart_style = 10


def slide_charts(prs: Presentation, results: list[Experiment]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    add_header(slide, "02", "품질과 속도 비교", "평균점수는 TC-01~16, 수행시간은 TC-01~18 실제 실행 사례의 중앙값")

    add_rect(slide, 0.48, 1.27, 6.1, 4.9, WHITE, line=LINE)
    add_rect(slide, 6.75, 1.27, 6.1, 4.9, WHITE, line=LINE)
    add_text(slide, 0.77, 1.48, 5.5, 0.38, "평균 품질점수", 16, NAVY, True)
    add_text(slide, 7.04, 1.48, 5.5, 0.38, "중앙 수행시간", 16, NAVY, True)

    score_data = ChartData()
    score_data.categories = [item.name for item in results]
    score_data.add_series("평균점수", [round(item.average, 1) for item in results])
    score_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.82), Inches(1.98), Inches(5.45), Inches(3.75), score_data
    ).chart
    _style_chart(score_chart, 100, 20, "0.0")
    score_chart.series[0].format.fill.solid()
    score_chart.series[0].format.fill.fore_color.rgb = BLUE

    time_data = ChartData()
    time_data.categories = [item.name for item in results]
    time_data.add_series("중앙시간", [round(item.median_seconds, 2) for item in results])
    time_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.09), Inches(1.98), Inches(5.45), Inches(3.75), time_data
    ).chart
    _style_chart(time_chart, 80, 20, "0.0초")
    time_chart.series[0].format.fill.solid()
    time_chart.series[0].format.fill.fore_color.rgb = GOLD

    fastest = min(results, key=lambda x: x.median_seconds)
    best = max(results, key=lambda x: x.average)
    add_text(slide, 0.75, 6.35, 5.75, 0.48,
             f"품질: {best.name}가 {best.average:.1f}점으로 가장 높음",
             12.5, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, 6.88, 6.35, 5.75, 0.48,
             f"속도: {fastest.name}가 {fastest.median_seconds:.2f}초로 가장 빠름",
             12.5, GOLD, True, PP_ALIGN.CENTER)


def slide_conclusion(prs: Presentation, results: list[Experiment]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_rect(slide, 0, 0, 13.333, 0.13, GOLD, radius=False)
    add_text(slide, 0.65, 0.42, 12.0, 0.56, "의견 및 결론", 27, WHITE, True)
    add_text(slide, 0.66, 0.98, 12.0, 0.36,
             "이번 1회 교차검증 실행에서 확인한 선택 기준", 11, RGBColor(177, 193, 218))

    best = max(results, key=lambda x: x.average)
    fastest = min(results, key=lambda x: x.median_seconds)
    points = [
        ("01", "품질 우선", f"{best.name} 조합이 평균 {best.average:.1f}점으로 가장 높아 최종 품질 확인에 유리했습니다."),
        ("02", "속도 우선", f"{fastest.name} 조합은 중앙 {fastest.median_seconds:.2f}초로 가장 빨라 반복 점검에 적합했습니다."),
        ("03", "해석 주의", "모호한 질문과 제한된 VOC 근거에서는 점수 편차가 커 모델 성능만으로 결과를 설명하기 어렵습니다."),
        ("04", "활용 방향", "빠른 반복 검증은 C, 서로 다른 모델의 교차 확인이 필요한 최종 검증은 A를 우선 활용하는 구성이 적절합니다."),
    ]
    y = 1.63
    accents = [BLUE, MINT, GOLD, SKY]
    for (number, label, body), accent in zip(points, accents):
        add_rect(slide, 0.70, y, 11.95, 0.96, RGBColor(31, 47, 78), line=RGBColor(56, 75, 108))
        add_text(slide, 0.90, y + 0.19, 0.48, 0.48, number, 12, WHITE, True, PP_ALIGN.CENTER)
        badge = slide.shapes[-1]
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.color.rgb = accent
        add_text(slide, 1.58, y + 0.12, 1.45, 0.30, label, 12, accent, True)
        add_text(slide, 1.58, y + 0.40, 10.60, 0.38, body, 12.5, WHITE)
        y += 1.08

    add_rect(slide, 0.70, 6.18, 11.95, 0.72, RGBColor(245, 183, 65), line=GOLD)
    add_text(slide, 0.95, 6.30, 11.45, 0.42,
             "결론  |  특정 모델의 절대 우위보다 목적에 맞는 조합 선택과 반복 교차검증이 중요합니다.",
             15, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, 0.72, 7.10, 11.9, 0.20, "※ 2026-07-16 실행 결과 · API 실패 N/A 0건", 8.5,
             RGBColor(153, 172, 202), False, PP_ALIGN.RIGHT)


def build_presentation() -> Path:
    results = load_results()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_summary(prs, results)
    slide_charts(prs, results)
    slide_conclusion(prs, results)
    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    print(build_presentation())
